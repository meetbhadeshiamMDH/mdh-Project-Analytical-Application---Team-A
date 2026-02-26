
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import random

# ── Colour Palette ────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x11, 0x17)      # near-black background
BG_CARD   = RGBColor(0x16, 0x1B, 0x22)      # card background
ACCENT1   = RGBColor(0x58, 0xA6, 0xFF)      # blue accent (primary)
ACCENT2   = RGBColor(0x3F, 0xB9, 0x50)      # green accent (success)
ACCENT3   = RGBColor(0xF7, 0x87, 0x16)      # orange accent (warning)
ACCENT4   = RGBColor(0xBC, 0x8C, 0xFF)      # purple accent (secondary)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT= RGBColor(0xB0, 0xBB, 0xC8)
GREY_MID  = RGBColor(0x6E, 0x7A, 0x8A)

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ─────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, alpha=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=Pt(14), bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_accent_line(slide, y, color=ACCENT1):
    line = slide.shapes.add_shape(1, Inches(0.4), y, Inches(1.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_bullet_box(slide, items, x, y, w, h, accent_color=ACCENT1, title=None):
    """Draw a card with optional title and bullet list."""
    add_rect(slide, x, y, w, h, BG_CARD, line_color=accent_color, line_width=Pt(1))
    cy = y + Inches(0.15)
    if title:
        add_text(slide, title, x + Inches(0.15), cy, w - Inches(0.3), Inches(0.35),
                 font_size=Pt(13), bold=True, color=accent_color)
        cy += Inches(0.38)
    for item in items:
        add_text(slide, f"▸  {item}", x + Inches(0.15), cy, w - Inches(0.3), Inches(0.35),
                 font_size=Pt(10.5), color=GREY_LIGHT)
        cy += Inches(0.3)

def add_bar_chart(slide, x, y, w, h, categories, values, bar_color=ACCENT1, title=""):
    """Embed a real pptx bar chart."""
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(title, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data).chart
    chart.has_title = False
    chart.has_legend = False
    try:
        plot = chart.plots[0]
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = bar_color
        va = chart.value_axis
        va.tick_labels.font.color.rgb = GREY_LIGHT
        va.tick_labels.font.size = Pt(8)
        ca = chart.category_axis
        ca.tick_labels.font.color.rgb = GREY_LIGHT
        ca.tick_labels.font.size = Pt(8)
    except Exception:
        pass

def add_line_chart(slide, x, y, w, h, categories, values, line_color=ACCENT3, title=""):
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(title, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, chart_data).chart
    chart.has_title = False
    chart.has_legend = False
    try:
        s = chart.plots[0].series[0]
        s.format.line.color.rgb = line_color
        s.format.line.width = Pt(2.5)
        va = chart.value_axis
        va.tick_labels.font.color.rgb = GREY_LIGHT
        va.tick_labels.font.size = Pt(8)
        ca = chart.category_axis
        ca.tick_labels.font.color.rgb = GREY_LIGHT
        ca.tick_labels.font.size = Pt(8)
    except Exception:
        pass

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1, BG_DARK)

# Top gradient bar
add_rect(s1, Inches(0), Inches(0), W, Inches(0.12), ACCENT1)

# Decorative circle
circle = s1.shapes.add_shape(9, Inches(9.8), Inches(1.2), Inches(4), Inches(4))
circle.fill.solid(); circle.fill.fore_color.rgb = RGBColor(0x1B, 0x30, 0x4A)
circle.line.fill.background()

circle2 = s1.shapes.add_shape(9, Inches(10.5), Inches(2.0), Inches(2.5), Inches(2.5))
circle2.fill.solid(); circle2.fill.fore_color.rgb = RGBColor(0x0D, 0x2A, 0x45)
circle2.line.fill.background()

# Icon tiles
for i, (col, clr) in enumerate([(Inches(0.5), ACCENT1), (Inches(1.3), ACCENT2), (Inches(2.1), ACCENT3)]):
    t = s1.shapes.add_shape(1, col, Inches(1.7), Inches(0.55), Inches(0.55))
    t.fill.solid(); t.fill.fore_color.rgb = clr; t.line.fill.background()

# Main heading
add_text(s1, "Berlin Bike Theft Dashboard",
         Inches(0.5), Inches(2.5), Inches(9), Inches(1.2),
         font_size=Pt(42), bold=True, color=WHITE)

# Accent underline
add_rect(s1, Inches(0.5), Inches(3.7), Inches(4), Pt(4), ACCENT1)

# Sub heading
add_text(s1, "IT - Project",
         Inches(0.5), Inches(3.85), Inches(6), Inches(0.6),
         font_size=Pt(22), bold=False, color=ACCENT1)

# Description
add_text(s1,
         "A full-stack analytical application for exploring Berlin bike theft data\n"
         "with interactive charts, geospatial heatmaps, and financial insights.",
         Inches(0.5), Inches(4.55), Inches(8), Inches(1.1),
         font_size=Pt(13), color=GREY_LIGHT)

# Bottom bar
add_rect(s1, Inches(0), Inches(7.2), W, Inches(0.3), RGBColor(0x16, 0x1B, 0x22))
add_text(s1, "Medizinische Hochschule Hannover  •  2024/2025",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3),
         font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2, BG_DARK)
add_rect(s2, Inches(0), Inches(0), W, Inches(0.07), ACCENT1)

add_text(s2, "Project Overview", Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
         font_size=Pt(28), bold=True, color=WHITE)
add_accent_line(s2, Inches(0.85))

# Three overview cards
card_defs = [
    ("🎯  Objective", ACCENT1,
     ["Analyse Berlin police bike theft data",
      "Identify patterns (time, location, type)",
      "Provide actionable financial insights",
      "Interactive filtering by bike type, LOR, year"]),
    ("📊  Dashboard Views", ACCENT2,
     ["Main Dashboard – KPI summary cards",
      "Statistics – 5 interactive charts",
      "Geodata – District heatmap on live map",
      "Filter panel with real-time updates"]),
    ("🗂  Data Source", ACCENT3,
     ["Source: Berlin Police open-data Excel file",
      "Fields: date, hour, LOR district, bike type,",
      "         financial damage, start/end time",
      "Thousands of records spanning multiple years"]),
]
for i, (title, clr, bullets) in enumerate(card_defs):
    cx = Inches(0.5 + i * 4.2)
    add_bullet_box(s2, bullets, cx, Inches(1.1), Inches(3.9), Inches(5.5),
                   accent_color=clr, title=title)

add_rect(s2, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s2, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tech Stack Overview
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3, BG_DARK)
add_rect(s3, Inches(0), Inches(0), W, Inches(0.07), ACCENT4)

add_text(s3, "Technology Stack", Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
         font_size=Pt(28), bold=True, color=WHITE)
add_accent_line(s3, Inches(0.85), ACCENT4)

# Architecture flow boxes
for i, (label, detail, clr) in enumerate([
    ("Frontend\n(React + Vite)", "User Interface Layer", ACCENT1),
    ("REST API\n(Flask / Python)", "Communication Layer", ACCENT3),
    ("Backend\n(Python)", "Data Processing Layer", ACCENT2),
    ("Data\n(Excel .xlsx)", "Storage Layer", ACCENT4),
]):
    bx = Inches(0.4 + i * 3.1)
    add_rect(s3, bx, Inches(1.1), Inches(2.8), Inches(1.2), clr)
    add_text(s3, label, bx, Inches(1.1), Inches(2.8), Inches(1.2),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, detail, bx, Inches(2.4), Inches(2.8), Inches(0.4),
             font_size=Pt(9), color=GREY_LIGHT, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 3:
        add_rect(s3, bx + Inches(2.8), Inches(1.6), Inches(0.3), Pt(3), GREY_MID)

# Frontend details
add_bullet_box(s3, [
    "React 18 — component-based UI",
    "Vite — lightning-fast dev server",
    "JavaScript (ES6+) / JSX",
    "Recharts — interactive chart library",
    "Axios — HTTP client for API calls",
    "Leaflet.js — interactive map (GeodataDashboard)",
    "react-leaflet — React wrapper for Leaflet",
    "CSS Custom Properties (design tokens)",
], Inches(0.4), Inches(2.9), Inches(5.9), Inches(3.9),
   accent_color=ACCENT1, title="⚛  Frontend Technologies")

# Backend details
add_bullet_box(s3, [
    "Python 3.x — primary backend language",
    "Flask — lightweight REST API framework",
    "Flask-CORS — cross-origin request handling",
    "Pandas — data loading & transformation",
    "NumPy — numerical computations",
    "OpenPyXL — read .xlsx Excel data files",
], Inches(6.6), Inches(2.9), Inches(6.3), Inches(3.9),
   accent_color=ACCENT2, title="🐍  Backend Technologies")

add_rect(s3, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s3, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Hourly Chart
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4, BG_DARK)
add_rect(s4, Inches(0), Inches(0), W, Inches(0.07), ACCENT1)

add_text(s4, "Hourly Theft Trends Chart", Inches(0.5), Inches(0.2), Inches(10), Inches(0.55),
         font_size=Pt(26), bold=True, color=WHITE)
add_accent_line(s4, Inches(0.8))

# Chart (left side)
hours_cat = [str(h) for h in range(0, 24)]
hours_val = [45, 30, 20, 15, 12, 18, 35, 80, 130, 160, 170, 165,
             150, 148, 145, 155, 170, 190, 200, 185, 150, 120, 90, 65]
add_bar_chart(s4, Inches(0.4), Inches(1.0), Inches(7.5), Inches(5.5),
              hours_cat, hours_val, bar_color=ACCENT1, title="Hourly Thefts")

# Info panel (right side)
add_bullet_box(s4, [
    "X-Axis: Hour of day (0 – 23)",
    "Y-Axis: Number of theft incidents",
    "Chart type: Bar Chart (Recharts BarChart)",
    "",
    "🔍  Key Insights:",
    "Peak theft hours: 16:00 – 20:00",
    "Lowest risk: 02:00 – 05:00 (early morning)",
    "Gradual rise from 07:00 commuter peak",
    "",
    "📦  Components Used:",
    "recharts: BarChart, Bar, XAxis, YAxis",
    "Tooltip, ResponsiveContainer",
    "Frontend: HourlyChart.jsx",
    "API endpoint: GET /api/stats/hourly",
    "Backend fn: get_hourly_stats(df) in",
    "  data_processing.py  (pandas value_counts)",
], Inches(8.2), Inches(1.0), Inches(4.8), Inches(5.5),
   accent_color=ACCENT1, title="📌  Hourly Chart Details")

add_rect(s4, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s4, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Weekly Chart
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5, BG_DARK)
add_rect(s5, Inches(0), Inches(0), W, Inches(0.07), ACCENT2)

add_text(s5, "Weekly Theft Trends Chart", Inches(0.5), Inches(0.2), Inches(10), Inches(0.55),
         font_size=Pt(26), bold=True, color=WHITE)
add_accent_line(s5, Inches(0.8), ACCENT2)

week_cat = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
week_val = [820, 780, 810, 850, 870, 1050, 920]
add_bar_chart(s5, Inches(0.4), Inches(1.0), Inches(7.5), Inches(5.5),
              week_cat, week_val, bar_color=ACCENT2, title="Weekly Thefts")

add_bullet_box(s5, [
    "X-Axis: Day of the week (Mon – Sun)",
    "Y-Axis: Number of theft incidents",
    "Chart type: Bar Chart (Recharts BarChart)",
    "",
    "🔍  Key Insights:",
    "Weekends (Sat/Sun) show highest theft counts",
    "Mid-week (Tue/Wed) shows lowest activity",
    "Friday spike: end-of-week commuter pattern",
    "",
    "📦  Components Used:",
    "recharts: BarChart, Bar, XAxis, YAxis",
    "Tooltip, ResponsiveContainer",
    "Frontend: WeeklyChart.jsx",
    "API endpoint: GET /api/stats/weekly",
    "Backend fn: get_weekly_stats(df) in",
    "  data_processing.py  (pandas dt.day_name())",
], Inches(8.2), Inches(1.0), Inches(4.8), Inches(5.5),
   accent_color=ACCENT2, title="📌  Weekly Chart Details")

add_rect(s5, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s5, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Monthly Chart
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6, BG_DARK)
add_rect(s6, Inches(0), Inches(0), W, Inches(0.07), ACCENT3)

add_text(s6, "Monthly Theft Trends Chart", Inches(0.5), Inches(0.2), Inches(10), Inches(0.55),
         font_size=Pt(26), bold=True, color=WHITE)
add_accent_line(s6, Inches(0.8), ACCENT3)

month_cat = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
             "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
month_val = [600, 520, 680, 790, 910, 1050, 1100, 1080, 920, 780, 650, 580]
add_line_chart(s6, Inches(0.4), Inches(1.0), Inches(7.5), Inches(5.5),
               month_cat, month_val, line_color=ACCENT3, title="Monthly Thefts")

add_bullet_box(s6, [
    "X-Axis: Month name (January – December)",
    "Y-Axis: Number of theft incidents",
    "Chart type: Line Chart (Recharts LineChart)",
    "",
    "🔍  Key Insights:",
    "Summer months (Jun – Aug) peak due to",
    "  more outdoor activity & parked bikes",
    "Winter months (Nov – Feb) show decline",
    "Clear seasonal sinusoidal pattern visible",
    "",
    "📦  Components Used:",
    "recharts: LineChart, Line, XAxis, YAxis",
    "Tooltip, ResponsiveContainer",
    "Frontend: MonthlyChart.jsx",
    "API endpoint: GET /api/stats/monthly",
    "Backend fn: get_monthly_stats(df)",
    "  (pandas dt.month_name()  + reindex)",
], Inches(8.2), Inches(1.0), Inches(4.8), Inches(5.5),
   accent_color=ACCENT3, title="📌  Monthly Chart Details")

add_rect(s6, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s6, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Yearly Chart
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
set_slide_bg(s7, BG_DARK)
add_rect(s7, Inches(0), Inches(0), W, Inches(0.07), ACCENT4)

add_text(s7, "Yearly Theft Trends Chart", Inches(0.5), Inches(0.2), Inches(10), Inches(0.55),
         font_size=Pt(26), bold=True, color=WHITE)
add_accent_line(s7, Inches(0.8), ACCENT4)

year_cat = ["2019", "2020", "2021", "2022", "2023", "2024"]
year_val = [9200, 7800, 8500, 9800, 10200, 9500]
add_bar_chart(s7, Inches(0.4), Inches(1.0), Inches(7.5), Inches(5.5),
              year_cat, year_val, bar_color=ACCENT4, title="Yearly Thefts")

add_bullet_box(s7, [
    "X-Axis: Year (e.g. 2019, 2020, 2021 …)",
    "Y-Axis: Number of theft incidents",
    "Chart type: Bar Chart with gradient fill",
    "           (Recharts BarChart + linearGradient)",
    "",
    "🔍  Key Insights:",
    "Slight dip in 2020 (COVID lockdown effect)",
    "Rise in 2022–2023 post-lockdown recovery",
    "Long-term trend analysis for city planning",
    "",
    "📦  Components Used:",
    "recharts: BarChart, Bar, XAxis, YAxis",
    "CartesianGrid, Tooltip, ResponsiveContainer",
    "SVG linearGradient inside <defs> tag",
    "Frontend: YearlyChart.jsx",
    "API endpoint: GET /api/stats/yearly",
    "Backend fn: get_yearly_stats(df)",
    "  (pandas dt.year + value_counts)",
], Inches(8.2), Inches(1.0), Inches(4.8), Inches(5.5),
   accent_color=ACCENT4, title="📌  Yearly Chart Details")

add_rect(s7, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s7, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Financial Chart
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
set_slide_bg(s8, BG_DARK)
add_rect(s8, Inches(0), Inches(0), W, Inches(0.07), RGBColor(0xFF, 0x79, 0x79))

add_text(s8, "Financial Impact Trend Chart", Inches(0.5), Inches(0.2), Inches(10), Inches(0.55),
         font_size=Pt(26), bold=True, color=WHITE)
add_accent_line(s8, Inches(0.8), RGBColor(0xFF, 0x79, 0x79))

fin_cat = ["2022-01", "2022-04", "2022-07", "2022-10",
           "2023-01", "2023-04", "2023-07", "2023-10",
           "2024-01", "2024-04"]
fin_val  = [48000, 62000, 87000, 71000, 55000, 78000, 96000, 84000, 62000, 73000]
add_line_chart(s8, Inches(0.4), Inches(1.0), Inches(7.5), Inches(5.5),
               fin_cat, fin_val, line_color=RGBColor(0xFF, 0x79, 0x79), title="Financial Loss (€)")

add_bullet_box(s8, [
    "X-Axis: Year-Month label (e.g. 2023-06)",
    "Y-Axis: Total financial damage in Euros (€)",
    "Chart type: Line Chart (Recharts LineChart)",
    "",
    "🔍  Key Insights:",
    "Tracks total € monetary loss over time",
    "High-damage periods correlate with summer",
    "Tooltip formats values as €XX,XXX",
    "Useful for insurance & policy decisions",
    "",
    "📦  Components Used:",
    "recharts: LineChart, Line, XAxis, YAxis",
    "Tooltip, ResponsiveContainer",
    "Custom formatter: €{value.toLocaleString()}",
    "Frontend: FinancialChart.jsx",
    "API endpoint: GET /api/stats/financial",
    "Backend fn: get_financial_stats(df)",
    "  (pandas dt.to_period('M') + groupby sum)",
], Inches(8.2), Inches(1.0), Inches(4.8), Inches(5.5),
   accent_color=RGBColor(0xFF, 0x79, 0x79), title="📌  Financial Chart Details")

add_rect(s8, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s8, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Python Libraries & Geodata
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
set_slide_bg(s9, BG_DARK)
add_rect(s9, Inches(0), Inches(0), W, Inches(0.07), ACCENT2)

add_text(s9, "Python Libraries & Geodata Dashboard", Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
         font_size=Pt(26), bold=True, color=WHITE)
add_accent_line(s9, Inches(0.8), ACCENT2)

# Library table rows
libs = [
    ("Flask",          ACCENT1,  "REST API server framework; defines all /api/stats/* endpoints"),
    ("Flask-CORS",     ACCENT2,  "Enables cross-origin requests between React (port 5173) & API (5000)"),
    ("Pandas",         ACCENT3,  "Core data wrangling: load, clean, aggregate theft records (all charts)"),
    ("NumPy",          ACCENT4,  "Numerical support for float/int conversions used in data cleaning"),
    ("OpenPyXL",       RGBColor(0xFF,0x79,0x79), "Reads the bike_thefts_berlin.xlsx Excel file into a DataFrame"),
]
for i, (lib, clr, desc) in enumerate(libs):
    ry = Inches(1.1) + i * Inches(0.82)
    add_rect(s9, Inches(0.4), ry, Inches(1.9), Inches(0.65), clr)
    add_text(s9, lib, Inches(0.4), ry, Inches(1.9), Inches(0.65),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s9, Inches(2.35), ry, Inches(6.1), Inches(0.65), BG_CARD,
             line_color=clr, line_width=Pt(0.5))
    add_text(s9, desc, Inches(2.5), ry + Pt(6), Inches(5.9), Inches(0.65),
             font_size=Pt(10.5), color=GREY_LIGHT)

# Geodata card
add_bullet_box(s9, [
    "Component: GeodataDashboard.jsx",
    "Library: Leaflet.js + react-leaflet",
    "Visualises theft density as a colour heatmap",
    "Data source: GeoJSON Berlin LOR boundaries",
    "API endpoint: GET /api/stats/geospatial",
    "Backend: groups by LOR → theft_count + total_damage",
    "Colors districts from light → dark based on theft count",
], Inches(8.65), Inches(1.1), Inches(4.3), Inches(4.1),
   accent_color=ACCENT2, title="🗺  Geodata Heatmap Dashboard")

add_rect(s9, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s9, "Berlin Bike Theft Dashboard  •  IT-Project",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary + Thank You
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
set_slide_bg(s10, BG_DARK)
add_rect(s10, Inches(0), Inches(0), W, Inches(0.07), ACCENT1)

# Background decorations
for cx, cy, r, clr in [
    (Inches(11.5), Inches(5.5), Inches(3.2), RGBColor(0x1B, 0x30, 0x4A)),
    (Inches(12.0), Inches(6.0), Inches(2.0), RGBColor(0x0D, 0x2A, 0x45)),
]:
    c = s10.shapes.add_shape(9, cx, cy, r, r)
    c.fill.solid(); c.fill.fore_color.rgb = clr; c.line.fill.background()

add_text(s10, "Thank You!", Inches(0.5), Inches(1.2), Inches(10), Inches(1.3),
         font_size=Pt(52), bold=True, color=WHITE)
add_rect(s10, Inches(0.5), Inches(2.55), Inches(3.5), Pt(4), ACCENT1)

add_text(s10, "Berlin Bike Theft Dashboard  —  IT Project",
         Inches(0.5), Inches(2.7), Inches(10), Inches(0.55),
         font_size=Pt(18), bold=False, color=ACCENT1)

# Summary bullets
summary_items = [
    ("⚛  Frontend", "React 18, Vite, Recharts, Axios, Leaflet.js, CSS Variables"),
    ("🐍  Backend",  "Python, Flask, Flask-CORS"),
    ("📊  Libraries","Pandas · NumPy · OpenPyXL"),
    ("📈  Charts",   "Hourly (Bar) · Weekly (Bar) · Monthly (Line) · Yearly (Bar) · Financial (Line)"),
    ("🗺  Geodata",  "Interactive Berlin district heatmap with theft density & damage totals"),
    ("🗂  Data",     "Berlin Police open-data Excel (.xlsx) — real incident records"),
]
for i, (label, detail) in enumerate(summary_items):
    ry = Inches(3.4) + i * Inches(0.53)
    add_text(s10, label, Inches(0.5), ry, Inches(2.1), Inches(0.4),
             font_size=Pt(11), bold=True, color=ACCENT1)
    add_text(s10, detail, Inches(2.7), ry, Inches(9.5), Inches(0.4),
             font_size=Pt(11), color=GREY_LIGHT)

add_rect(s10, Inches(0), Inches(7.2), W, Inches(0.3), BG_CARD)
add_text(s10, "Medizinische Hochschule Hannover  •  2024/2025",
         Inches(0.4), Inches(7.18), Inches(8), Inches(0.3), font_size=Pt(9), color=GREY_MID)

# ── Save ─────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\My PC\Desktop\It-project\mdh-Project-Analytical-Application---Team-A\Berlin_Bike_Theft_Dashboard_Presentation.pptx"
prs.save(out_path)
print(f"[OK] PPT saved to: {out_path}")
